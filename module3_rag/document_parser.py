import re
import sys
import tempfile
import os
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))
from config import SUPPORTED_EXTENSIONS


def _parse_pdf(file_path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(file_path)
    pages = []
    for page in reader.pages:
        # layout mode preserves table row order instead of reading column-by-column
        text = page.extract_text(extraction_mode="layout") or ""
        text = re.sub(r'[ \t]+', ' ', text)
        text = re.sub(r'\n{3,}', '\n\n', text)
        pages.append(text.strip())
    return "\n\n".join(p for p in pages if p)


def _parse_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _parse_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
    return "\n".join(texts)


def _parse_xlsx(file_path: str) -> str:
    import pandas as pd
    return pd.read_excel(file_path, engine="openpyxl").to_string(index=False)


def _parse_csv(file_path: str) -> str:
    import pandas as pd
    return pd.read_csv(file_path, encoding="utf-8", errors="replace").to_string(index=False)


_PARSERS: dict = {
    ".pdf":  _parse_pdf,
    ".docx": _parse_docx,
    ".txt":  _parse_txt,
    ".md":   _parse_txt,
    ".pptx": _parse_pptx,
    ".xlsx": _parse_xlsx,
    ".csv":  _parse_csv,
}


def parse_document(file_path: str) -> str:
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"Fichier introuvable : {file_path}")

    ext = path.suffix.lower()
    if ext not in _PARSERS:
        raise ValueError(f"Format '{ext}' non supporté. Formats acceptés : {list(_PARSERS.keys())}")

    text = _PARSERS[ext](file_path)
    if not text.strip():
        raise ValueError(f"Le document '{path.name}' semble vide ou illisible.")

    return text


if __name__ == "__main__":
    content = "Bonjour étudiant.\nCeci est un test.\nBonne chance !"
    with tempfile.NamedTemporaryFile(mode="w", suffix=".txt", delete=False, encoding="utf-8") as f:
        f.write(content)
        tmp = f.name

    assert parse_document(tmp).strip() == content.strip()
    os.unlink(tmp)

    try:
        parse_document("fichier.xyz")
    except ValueError:
        pass

    missing = str(Path(tempfile.gettempdir()) / "edubot_missing.txt")
    try:
        parse_document(missing)
    except FileNotFoundError:
        pass

    print("document_parser.py: OK")
